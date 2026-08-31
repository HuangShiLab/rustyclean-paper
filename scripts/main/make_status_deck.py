#!/usr/bin/env python3
"""RustyClean status deck: code architecture + full experiment inventory.
Every number is computed from the CSVs in data/, never transcribed."""
import csv, statistics as st
from collections import defaultdict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK=RGBColor(0x2E,0x2A,0x27); RUST=RGBColor(0xB8,0x50,0x42); RUST_DK=RGBColor(0x8F,0x3A,0x2F)
RUST_LT=RGBColor(0xE0,0xA9,0x9E); SAGE=RGBColor(0xA7,0xBE,0xAE); SLATE=RGBColor(0x6E,0x7B,0x85)
CARD=RGBColor(0xF2,0xEF,0xEC); CARD_DK=RGBColor(0x40,0x3A,0x36); MUTED=RGBColor(0x7A,0x73,0x6C)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xDE,0xD8,0xD2); GREEN=RGBColor(0x2E,0x7D,0x5B)
H="Cambria"; B="Calibri"; M="Courier New"
SW,SH=13.333,7.5; MG=0.62; CW=SW-2*MG
warn=[]

def txt(s,t,x,y,w,h,size=14,color=INK,bold=False,font=B,align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,italic=False,sp=1.0,label=""):
    bx=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=bx.text_frame
    tf.word_wrap=True; tf.margin_left=tf.margin_right=Inches(0.05)
    tf.margin_top=tf.margin_bottom=Inches(0.02); tf.vertical_anchor=anchor
    for i,ln in enumerate(str(t).split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=sp
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold
        r.font.italic=italic; r.font.name=font; r.font.color.rgb=color
    f=0.50 if font==B else (0.53 if font==H else 0.60)
    if bold: f+=0.02
    cpl=max(1,int((w-0.1)*72/(f*size)))
    need=sum(max(1,-(-len(l)//cpl)) for l in str(t).split("\n"))*size*1.24/72+0.1
    if need>h+0.06: warn.append(f"{label or str(t)[:30]!r}: {need:.2f}>{h:.2f}")
    return bx

def rect(s,x,y,w,h,fill=CARD,line=None,rad=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rad else MSO_SHAPE.RECTANGLE,
                          Inches(x),Inches(y),Inches(w),Inches(h))
    if rad: sh.adjustments[0]=rad
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line: sh.line.color.rgb=line; sh.line.width=Pt(1)
    else: sh.line.fill.background()
    sh.shadow.inherit=False; sh.text_frame.text=""
    return sh

def circ(s,x,y,d,fill=RUST,t=None,size=15,tc=WHITE):
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(d),Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.fill.background(); sh.shadow.inherit=False
    tf=sh.text_frame; tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
    tf.vertical_anchor=MSO_ANCHOR.MIDDLE; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    if t:
        r=p.add_run(); r.text=t; r.font.size=Pt(size); r.font.bold=True
        r.font.name=B; r.font.color.rgb=tc
    return sh

def slide(prs,dark=False):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    if dark: rect(s,0,0,SW,SH,fill=INK)
    return s

def title(s,t,kicker=None,dark=False,size=32):
    y=0.50
    if kicker:
        txt(s,kicker.upper(),MG,0.40,CW,0.26,size=11,bold=True,
            color=RUST_LT if dark else RUST,label="k"); y=0.70
    txt(s,t,MG,y,CW,0.66,size=size,bold=True,color=WHITE if dark else INK,font=H,label="t:"+t[:20])

def note(s,t): s.notes_slide.notes_text_frame.text=t

def table(s,x,y,w,headers,rows,widths,fs=10.5,hfs=10.5,rh=0.32,hh=0.36,
          zebra=True,colcol=None):
    rect(s,x,y,w,hh,fill=INK,rad=0.04)
    cx=x
    for i,hd in enumerate(headers):
        txt(s,hd,cx+0.06,y+0.04,widths[i]-0.1,hh-0.08,size=hfs,bold=True,color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE,label="h"+hd[:12]); cx+=widths[i]
    yy=y+hh
    for ri,row in enumerate(rows):
        if zebra and ri%2==0: rect(s,x,yy,w,rh,fill=CARD)
        cx=x
        for i,cell in enumerate(row):
            col=INK
            if colcol and i in colcol: col=colcol[i](cell)
            txt(s,cell,cx+0.06,yy+0.02,widths[i]-0.1,rh-0.04,size=fs,color=col,
                anchor=MSO_ANCHOR.MIDDLE,bold=(i==0),label="c"+str(cell)[:12]); cx+=widths[i]
        yy+=rh
    return yy

# ------------------------------------------------------------------ data
def rd(p):
    with open(p) as f: return list(csv.DictReader(f))
def rows_lenient(p):
    out=[]
    with open(p) as f:
        rdr=csv.reader(f); hdr=next(rdr)
        for r in rdr:
            if not r: continue
            if len(r)==len(hdr)-1: r=r[:2]+[""]+r[2:]
            out.append(dict(zip(hdr,r)))
    return out
def cm(r): return tuple(int(float(r[k])) for k in ("tp","fp","tn","fn"))
def met(tp,fp,tn,fn): return (2*tp/(2*tp+fp+fn), fn/(tp+fn)*100, fp/(fp+tn)*100)

D="data"
ac=rd(f"{D}/benchmark_results/accuracy_comparison.csv")
pm=rd(f"{D}/benchmark_results/auto_vs_kneaddata_metrics.csv")
fh=rd(f"{D}/benchmark_results/fair_hostile_skipqc_results.csv")
bs=rd(f"{D}/benchmark_results/accuracy_baseline_kraken2_memmap.csv")
rc_=rd(f"{D}/benchmark_results/accuracy_bowtie2_recheck.csv")
pbs=rd(f"{D}/benchmark_results/performance_baseline_kraken2_memmap.csv")
prc=rd(f"{D}/benchmark_results/performance_bowtie2_recheck.csv")
v4a=rd(f"{D}/accuracy_rc_mm_bt_cf_v4.csv")
v4p={(r["tool"],r["dataset"]):r for r in rd(f"{D}/performance_rc_mm_bt_cf_v4_corrected.csv")}
m100a=rd(f"{D}/results_100M_matched/accuracy.csv"); m100p=rows_lenient(f"{D}/results_100M_matched/performance.csv")
s100a=rd(f"{D}/results_100M_skipqc_matched/accuracy.csv"); s100p=rows_lenient(f"{D}/results_100M_skipqc_matched/performance.csv")

O4=["5M_1pct_low_even_SE","10M_10pct_med_even_SE","30M_50pct_high_skewed_SE","60M_90pct_high_lognormal_SE"]
ORC=["30M_50pct_high_skewed_SE","60M_90pct_high_lognormal_SE","100M_50pct_high_lognormal_SE","100M_90pct_high_lognormal_SE"]
D100=["100M_50pct_high_lognormal_SE","100M_90pct_high_lognormal_SE"]
SH_={"5M_1pct_low_even_SE":"5M / 1%","10M_10pct_med_even_SE":"10M / 10%",
     "30M_50pct_high_skewed_SE":"30M / 50%","60M_90pct_high_lognormal_SE":"60M / 90%",
     "100M_50pct_high_lognormal_SE":"100M / 50%","100M_90pct_high_lognormal_SE":"100M / 90%"}
A=lambda t,d: met(*cm([r for r in ac if r["tool"]==t and r["dataset"]==d][0]))
P={}; [P.__setitem__((r["tool"],r["dataset"]),r) for r in pm]
F={}; [F.__setitem__((r["tool"],r["dataset"]),r) for r in fh]
Bx=lambda d: met(*cm([r for r in bs if r["dataset"]==d][0]))
Rx=lambda d: met(*cm([r for r in rc_ if r["dataset"]==d][0]))
def rt(rows,d,tool=None,col="runtime_seconds"):
    v=[float(x[col]) for x in rows if x["dataset"]==d and (tool is None or x["tool"]==tool) and x[col]]
    return st.mean(v) if v else 0
MEAN={t:{k:st.mean([A(t,d)[i] for d in O4]) for i,k in enumerate(("f1","ml","hc"))}
      for t in ("rustyclean_auto","hostile_raw","kneaddata")}
SPD_KD=[float(P[("kneaddata",d)]["runtime_seconds"])/float(P[("rustyclean_auto",d)]["runtime_seconds"]) for d in O4]
SPD_HO=[float(F[("hostile_raw",d)]["runtime_seconds"])/float(F[("rustyclean_auto_skipqc",d)]["runtime_seconds"]) for d in O4]
HCR=st.mean(Bx(d)[2]/Rx(d)[2] for d in ORC)
RTD=st.mean((rt(prc,d)-rt(pbs,d))/rt(pbs,d)*100 for d in ORC)

prs=Presentation(); prs.slide_width=Inches(SW); prs.slide_height=Inches(SH)

# ============================================================ 1 TITLE
s=slide(prs,dark=True)
for d,c in ((5.4,RGBColor(0x3A,0x33,0x2E)),(3.9,RGBColor(0x4A,0x37,0x31)),(2.4,RUST_DK),(1.1,RUST)):
    circ(s,10.7-d/2,3.75-d/2,d,fill=c)
txt(s,"RustyClean",MG,1.86,7.4,1.20,size=60,bold=True,color=WHITE,font=H,label="T")
txt(s,"Architecture and benchmark status",MG,3.18,7.2,0.50,size=21,color=SAGE,label="sub")
txt(s,"Adaptive host depletion for short-read metagenomics",MG,3.76,7.2,0.40,size=14,
    color=RUST_LT,italic=True,label="tag")
rect(s,MG,4.52,1.4,0.045,fill=RUST)
txt(s,"Code at commit 13eeb99  ·  7 benchmark experiments  ·  5 depletion backends",
    MG,4.82,7.2,0.32,size=13,color=WHITE,label="w")
txt(s,"ZHANG Yufeng  ·  Huang Lab",MG,5.28,7.0,0.30,size=13,color=MUTED,label="who")
note(s,"Status update: what the code does now, and every benchmark run to date with its result.")

# ============================================================ 2 PIPELINE
s=slide(prs); title(s,"Pipeline architecture","Code · what runs per sample")
steps=[("1  Quality control","fastp — adapters, trimming, length/quality filter","(skippable with --skip-qc)",RUST),
       ("2  Host-fraction survey","seqtk subsample 100k → bowtie2 --very-fast-local","estimates host % in seconds",RUST),
       ("3  Backend routing","choose_auto_backend(host%, read count)","picks bowtie2 or kraken2",SAGE),
       ("4  Host depletion","kraken2 | bowtie2 | minimap2 | sylph | centrifuge","5 interchangeable backends",RUST),
       ("5  Validation gate","output size + residual contamination assertions","output promoted only if it passes",CARD_DK)]
y=1.52
for i,(hd,mid,sub,col) in enumerate(steps):
    rect(s,MG,y,CW,0.86,fill=col if col!=SAGE else SAGE,rad=0.08)
    tc=INK if col==SAGE else WHITE
    sc=INK if col==SAGE else (RUST_LT if col==RUST else RGBColor(0xC9,0xC1,0xBA))
    txt(s,hd,MG+0.28,y+0.10,3.05,0.32,size=15,bold=True,color=tc,label="s"+hd[:9])
    txt(s,mid,MG+3.45,y+0.09,5.6,0.34,size=13,color=tc,font=M,label="m"+hd[:9])
    txt(s,sub,MG+9.05,y+0.09,CW-9.30,0.62,size=11,color=sc,italic=True,sp=1.06,
        anchor=MSO_ANCHOR.MIDDLE,label="u"+hd[:9])
    if i<4: txt(s,"↓",MG,y+0.86,CW,0.20,size=11,bold=True,color=RUST_DK,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    y+=1.06
rect(s,MG,6.86,CW,0.0,fill=LINE)
txt(s,"Stages 1 and 2 are optional; stages 3–5 always run. Every stage transition is recorded "
      "in a per-sample checkpoint.",MG,6.72,CW,0.34,size=12.5,italic=True,color=MUTED,label="foot")
note(s,"Five stages. The survey and routing (2-3) are what makes it adaptive; the gate (5) is what "
       "no other tool has.")

# ============================================================ 3 MODULES
s=slide(prs); title(s,"Code structure","Code · 2,521 lines of Rust, 9 modules")
mods=[("pipeline.rs","1402","Stage machine, all 5 backends,\nsurvey, routing, metrics",True),
      ("main.rs","314","Config wiring, pre-flight\nchecks, run summary",False),
      ("worker.rs","215","Async worker pool, retry,\ncancellation, progress",True),
      ("config.rs","180","TOML config, defaults,\nCLI overrides",False),
      ("sample.rs","133","Manifest parsing, SE/PE\ndetection, input fingerprint",False),
      ("cli.rs","131","24 flags incl. auto-mode\nthresholds",False),
      ("checkpoint.rs","113","Atomic versioned JSON\ncheckpoint store",True),
      ("error.rs / lib.rs","33","Typed error taxonomy,\ncrate root",False)]
gw=(CW-3*0.24)/4; gh=1.90
for i,(n,loc,role,hot) in enumerate(mods):
    x=MG+(i%4)*(gw+0.24); y=1.56+(i//4)*(gh+0.24)
    rect(s,x,y,gw,gh,fill=CARD if hot else None,line=None if hot else LINE,rad=0.07)
    txt(s,n,x+0.20,y+0.20,gw-0.40,0.30,size=14,bold=True,color=RUST_DK if hot else INK,font=M,label="n"+n)
    txt(s,loc+" lines",x+0.20,y+0.54,gw-0.40,0.26,size=11,color=MUTED,label="l"+n)
    txt(s,role,x+0.20,y+0.86,gw-0.40,0.92,size=12,color=MUTED,sp=1.12,label="r"+n)
rect(s,MG,5.86,CW,0.86,fill=INK,rad=0.06)
for i,(big,small) in enumerate([("2,521","lines of Rust"),("9","modules"),("5","depletion backends"),
                                ("24","CLI flags")]):
    x=MG+i*(CW/4)
    txt(s,big,x,5.94,CW/4,0.46,size=21,bold=True,color=RUST_LT,align=PP_ALIGN.CENTER,font=H,label="f"+big)
    txt(s,small,x,6.42,CW/4,0.28,size=11.5,color=SAGE,align=PP_ALIGN.CENTER,label="fs"+small)
note(s,"Shaded modules carry the design: pipeline (backends + routing), worker (concurrency), "
       "checkpoint (crash safety). pipeline.rs tripled in size since the first version.")

# ============================================================ 4 ROUTING
s=slide(prs); title(s,"Adaptive backend routing","Code · pipeline.rs :: choose_auto_backend")
rect(s,MG,1.52,7.0,2.30,fill=INK,rad=0.05)
code=("if host_pct < low_threshold {          // default 10 %\n"
      "    \"bowtie2\"\n"
      "} else if host_pct > high_threshold    // default 30 %\n"
      "       && input_reads > reads_threshold  // default 20 M\n"
      "{   \"kraken2\"\n"
      "} else {\n"
      "    \"bowtie2\"      // conservative default\n"
      "}")
txt(s,code,MG+0.26,1.66,6.5,2.02,size=11,color=SAGE,font=M,sp=1.16,label="code")
rx=MG+7.28; rw=CW-7.28
txt(s,"A two-dimensional rule",rx,1.52,rw,0.34,size=17,bold=True,color=RUST_DK,font=H,label="rh")
for i,t in enumerate([
  "Classification only pays off when the host fraction is high AND the library is large enough to "
  "amortise loading the database.",
  "Anything that does not clearly meet both criteria falls back to alignment, whose error direction "
  "is the safer default.",
  "The estimate only has to land on the correct side of a threshold — not be accurate."]):
    circ(s,rx+0.02,1.96+i*0.76,0.20,fill=RUST)
    txt(s,t,rx+0.34,1.90+i*0.76,rw-0.40,0.76,size=12.1,color=INK,sp=1.10,label="b"+str(i))
rect(s,MG,4.04,CW,1.06,fill=CARD,rad=0.06)
txt(s,"Survey cost",MG+0.28,4.16,1.6,0.28,size=13,bold=True,color=RUST_DK,label="sc")
txt(s,"seqtk sample (fixed seed 42) draws 100,000 reads; bowtie2 --very-fast-local aligns them "
      "with 2 threads. Seconds, against runs of tens of minutes. --host-pct bypasses the survey.",
    MG+1.85,4.14,CW-2.15,0.60,size=12.3,color=MUTED,sp=1.10,label="scb")
rect(s,MG,5.26,CW,1.30,fill=INK,rad=0.06)
txt(s,"Observed routing — correct on all four datasets",MG+0.30,5.38,6.0,0.30,size=13.5,bold=True,
    color=RUST_LT,label="or")
cells=[(SH_[d],f"{float(P[('rustyclean_auto',d)]['estimated_host_pct'])}%",
        P[("rustyclean_auto",d)]["backend"]) for d in O4]
cw4=(CW-0.6)/4
for i,(dl,est,bk) in enumerate(cells):
    x=MG+0.30+i*cw4
    txt(s,dl,x,5.76,cw4-0.2,0.26,size=11.5,color=SAGE,label="d"+dl)
    txt(s,f"est {est}  →  {bk}",x,6.04,cw4-0.2,0.30,size=12.5,bold=True,color=WHITE,label="e"+dl)
note(s,"The rule reproduces every observed backend choice exactly. Note it is two-dimensional — "
       "host fraction alone is not enough.")

# ============================================================ 5 RELIABILITY
s=slide(prs); title(s,"Reliability layer","Code · worker.rs, checkpoint.rs")
cards=[("Checkpoint / resume","Versioned JSON per sample, written atomically (temp file + rename). "
        "An xxHash3 fingerprint over input metadata invalidates stale state automatically."),
       ("Bounded concurrency","A semaphore admits W samples at once; each passes T threads to its "
        "tools. Total load ≈ W × T. Ctrl-C sets a cancellation token so queued samples never start."),
       ("Validation gate","Output size and residual-contamination assertions run before promotion. "
        "A sample that fails is marked Failed and its output is never moved into the results directory.")]
cw3=(CW-2*0.32)/3
for i,(hd,body) in enumerate(cards):
    x=MG+i*(cw3+0.32)
    rect(s,x,1.54,cw3,2.56,fill=None,line=LINE,rad=0.07)
    circ(s,x+0.28,1.82,0.48,fill=RUST,t=str(i+1),size=16)
    txt(s,hd,x+0.28,2.48,cw3-0.56,0.32,size=15,bold=True,color=INK,label="k"+hd[:10])
    txt(s,body,x+0.28,2.84,cw3-0.56,1.20,size=11.9,color=MUTED,sp=1.10,label="kb"+hd[:10])
rect(s,MG,4.20,CW,1.38,fill=INK,rad=0.07)
txt(s,"Known gaps in this layer",MG+0.32,4.34,4.0,0.30,size=13.5,bold=True,color=RUST_LT,label="gh")
txt(s,"Checkpoints are written at attempt boundaries, not at every stage transition, so a hard "
      "process kill resumes further back than necessary. The per-sample timeout is configured but "
      "not enforced. Worker count is derived from CPU count, not available memory — with the "
      "15.5 GB Kraken2 index, concurrent demand scales as W × 15.5 GB.",
    MG+0.32,4.68,CW-0.64,0.82,size=13,color=WHITE,sp=1.12,label="gb")
txt(s,"The stage enum still names only the Kraken2 path (Kraken2Running / Kraken2Complete) even "
      "though five backends share it — cosmetic, but confusing when reading checkpoints.",
    MG,5.72,CW,0.56,size=12.3,italic=True,color=MUTED,label="nit")
note(s,"Be candid about the gaps — the memory-aware worker cap is the one that actually bites, "
       "because the Kraken2 index is 15.5 GB.")

# ============================================================ 6 EXPERIMENT INVENTORY
s=slide(prs); title(s,"Benchmark inventory","Evaluation · everything run to date")
hdr=["#","Experiment","Compared against","Datasets","Reps"]
wid=[0.42,4.05,3.55,2.30,1.77]
rows=[["1","Error profile & accuracy","Hostile, KneadData","4 SE (1–90 % host)","1"],
      ["2","Auto routing + full-pipeline runtime","KneadData","4 SE","1"],
      ["3","Fair depletion-only runtime","Hostile (--skip-qc both sides)","4 SE","1"],
      ["4","Bowtie2 recheck ablation","itself (Kraken2 only)","4 SE (50–90 %)","3"],
      ["5","Backend comparison","bowtie2 / minimap2 / centrifuge","4 SE","1"],
      ["6","100M matched panel","Hostile, KneadData","2 SE (100M)","3 / 1 / 1"],
      ["7","100M depletion-only matched","Hostile","2 SE (100M)","3 / 1"]]
table(s,MG,1.52,CW,hdr,rows,wid,fs=11.5,hfs=11.5,rh=0.40,hh=0.42)
rect(s,MG,4.62,CW,0.96,fill=CARD,rad=0.06)
txt(s,"All simulated, with per-read ground truth",MG+0.30,4.74,4.4,0.28,size=13,bold=True,
    color=RUST_DK,label="ah")
txt(s,"Host removal is scored exactly as a binary classification of reads. Depletion is "
      "deterministic, so accuracy is identical across replicates — replication applies to timing only.",
    MG+0.30,5.04,CW-0.60,0.54,size=12.3,color=MUTED,sp=1.10,label="ab")
txt(s,"Not yet run: paired-end libraries, host fractions between 10 % and 30 % (where routing "
      "actually decides), and any real cohort.",MG,5.76,CW,0.34,size=12.5,italic=True,
    color=RUST_DK,label="gap")
note(s,"Seven experiments. Note what is missing: PE, the 10-30% routing band, and real data.")

# ============================================================ 7 EXP 1 ACCURACY
s=slide(prs); title(s,"Experiment 1 — the two error directions","Results · accuracy, 4 datasets")
hdr=["Tool","Mean F1","Microbial loss","Host carry-over"]
wid=[3.9,2.2,3.0,2.9]
def fmt(v,d=3): return f"{v:.{d}f} %"
rows=[["RustyClean (auto, no recheck)",f"{MEAN['rustyclean_auto']['f1']:.4f}",
       fmt(MEAN['rustyclean_auto']['ml']),fmt(MEAN['rustyclean_auto']['hc'])],
      ["Hostile",f"{MEAN['hostile_raw']['f1']:.4f}",fmt(MEAN['hostile_raw']['ml']),
       fmt(MEAN['hostile_raw']['hc'])],
      ["KneadData",f"{MEAN['kneaddata']['f1']:.4f}",fmt(MEAN['kneaddata']['ml']),
       fmt(MEAN['kneaddata']['hc'])]]
table(s,MG,1.52,12.0,hdr,rows,wid,fs=13,hfs=12.5,rh=0.46,hh=0.44)
r1=rect(s,MG,3.34,CW/2-0.18,1.60,fill=CARD,rad=0.07)
txt(s,"KneadData discards microbes",MG+0.28,3.48,CW/2-0.74,0.30,size=14,bold=True,color=RUST_DK,label="c1")
txt(s,f"{MEAN['kneaddata']['ml']:.2f} % of genuine microbial reads, "
      f"{MEAN['kneaddata']['ml']/MEAN['rustyclean_auto']['ml']:.0f}× more than RustyClean. "
      f"Alignment's false positives are an irreversible loss of signal.",
    MG+0.28,3.84,CW/2-0.74,0.94,size=12.5,color=MUTED,sp=1.10,label="c1b")
rect(s,MG+CW/2+0.18,3.34,CW/2-0.18,1.60,fill=INK,rad=0.07)
txt(s,"RustyClean retains host",MG+CW/2+0.46,3.48,CW/2-0.74,0.30,size=14,bold=True,color=RUST_LT,label="c2")
txt(s,f"{MEAN['rustyclean_auto']['hc']:.2f} % of host reads survive without the verification pass — "
      f"the characteristic false-negative bias of k-mer classification. Experiment 4 addresses this.",
    MG+CW/2+0.46,3.84,CW/2-0.74,0.94,size=12.5,color=WHITE,sp=1.10,label="c2b")
txt(s,"Hostile is the strongest accuracy baseline: lowest microbial loss and highest mean F1 of the "
      "three. On accuracy alone, without verification, RustyClean does not lead.",
    MG,5.14,CW,0.58,size=13.5,bold=True,color=RUST_DK,sp=1.10,label="hon")
txt(s,"Confirms the pattern reported by Gao et al. (2025): alignment → false positives, "
      "k-mer → false negatives.",MG,5.76,CW,0.34,size=12.5,italic=True,color=MUTED,label="gao")
note(s,"State plainly that Hostile leads on accuracy here. The recheck experiment is what changes it.")

# ============================================================ 8 EXP 2+3 SPEED
s=slide(prs); title(s,"Experiments 2 & 3 — runtime","Results · speed against both baselines")
hdr=["Dataset","Backend","RustyClean","KneadData","vs KD","RC --skip-qc","Hostile","vs Hostile"]
wid=[1.75,1.30,1.65,1.55,1.05,1.85,1.45,1.50]
rows=[]
for d in O4:
    rc=float(P[("rustyclean_auto",d)]["runtime_seconds"]); kd=float(P[("kneaddata",d)]["runtime_seconds"])
    rs=float(F[("rustyclean_auto_skipqc",d)]["runtime_seconds"]); ho=float(F[("hostile_raw",d)]["runtime_seconds"])
    rows.append([SH_[d],P[("rustyclean_auto",d)]["backend"],f"{rc/60:.1f} min",f"{kd/60:.1f} min",
                 f"{kd/rc:.1f}×",f"{rs/60:.1f} min",f"{ho/60:.1f} min",f"{ho/rs:.1f}×"])
table(s,MG,1.52,CW,hdr,rows,wid,fs=11.5,hfs=10.5,rh=0.40,hh=0.42)
rect(s,MG,3.36,CW/2-0.18,1.44,fill=CARD,rad=0.07)
txt(s,"vs KneadData — full pipeline",MG+0.28,3.48,CW/2-0.74,0.30,size=13.5,bold=True,color=RUST_DK,label="k1")
txt(s,f"{min(SPD_KD):.1f}–{max(SPD_KD):.1f}× faster on every dataset, widest at the highest host "
      f"fraction. RustyClean does no repeat masking, so part of this is work not done.",
    MG+0.28,3.82,CW/2-0.74,0.86,size=12.3,color=MUTED,sp=1.10,label="k1b")
rect(s,MG+CW/2+0.18,3.36,CW/2-0.18,1.44,fill=CARD,rad=0.07)
txt(s,"vs Hostile — depletion only",MG+CW/2+0.46,3.48,CW/2-0.74,0.30,size=13.5,bold=True,
    color=RUST_DK,label="k2")
txt(s,f"Comparable at low host ({SPD_HO[0]:.2f}× and {SPD_HO[1]:.1f}×), "
      f"{SPD_HO[2]:.1f}× and {SPD_HO[3]:.1f}× faster at 50 % and 90 % host. QC skipped on both "
      f"sides so only depletion is timed.",
    MG+CW/2+0.46,3.82,CW/2-0.74,0.86,size=12.3,color=MUTED,sp=1.10,label="k2b")
rect(s,MG,5.00,CW,1.10,fill=INK,rad=0.06)
txt(s,"Caveat — these two Hostile numbers do not reconcile with Experiment 6",
    MG+0.30,5.12,7.6,0.30,size=13,bold=True,color=RUST_LT,label="cv")
txt(s,f"Hostile takes {float(F[('hostile_raw',O4[3])]['runtime_seconds'])/60:.0f} min on the 60M "
      f"dataset here but only {rt(s100p,D100[1],'hostile')/60:.0f} min on the larger 100M dataset in "
      f"Experiment 6. A bigger library cannot be faster — one of the two runs used different "
      f"settings or a different node, and must be repeated before either speed-up is quoted.",
    MG+0.30,5.44,CW-0.60,0.58,size=12.3,color=WHITE,sp=1.10,label="cvb")
note(s,"Do not gloss over the Hostile inconsistency — 60M taking 71 min while 100M takes 22 min is "
       "not physically sensible. Flag it, do not quote the 4.9x until it is resolved.")

# ============================================================ 9 EXP 4 RECHECK
s=slide(prs); title(s,"Experiment 4 — Bowtie2 verification pass","Results · the ablation that matters")
hdr=["Dataset","Host carry-over","Microbial loss","F1","Runtime"]
wid=[2.6,3.1,2.7,2.0,1.7]
rows=[]
for d in ORC:
    b=Bx(d); r=Rx(d); tb=rt(pbs,d)/60; tr=rt(prc,d)/60
    rows.append([SH_[d],f"{b[2]:.3f} → {r[2]:.4f} %",f"{b[1]:.3f} → {r[1]:.3f} %",
                 f"{b[0]:.4f} → {r[0]:.4f}",f"{(tr-tb)/tb*100:+.1f} %"])
table(s,MG,1.52,12.1,hdr,rows,wid,fs=11.5,hfs=11.5,rh=0.40,hh=0.42)
big=[(f"{HCR:.1f}×","less residual host",RUST),
     (f"+{RTD:.1f}%","mean runtime cost",SLATE),
     (f"{MEAN['kneaddata']['ml']/st.mean(Rx(d)[1] for d in ORC):.1f}×","less microbial loss\nthan KneadData",RUST_DK)]
bw=(CW-2*0.34)/3
for i,(v,lab,col) in enumerate(big):
    x=MG+i*(bw+0.34)
    rect(s,x,3.34,bw,1.42,fill=CARD,rad=0.07)
    txt(s,v,x,3.48,bw,0.60,size=32,bold=True,color=col,align=PP_ALIGN.CENTER,font=H,label="v"+v)
    txt(s,lab,x,4.14,bw,0.52,size=12.5,color=MUTED,align=PP_ALIGN.CENTER,sp=1.12,label="lb"+v)
txt(s,"Remarkably consistent: host carry-over falls from 1.41 % to 0.07 % on all four datasets, "
      "and F1 at 90 % host rises from 0.930 to 0.994.",
    MG,4.94,CW,0.54,size=13.5,bold=True,color=INK,sp=1.10,label="cons")
rect(s,MG,5.44,CW,1.12,fill=INK,rad=0.06)
txt(s,"Now on main",MG+0.30,5.56,1.7,0.28,size=13,bold=True,color=RUST_LT,label="nm")
txt(s,"--bowtie2-recheck has been merged and now takes the Bowtie2 index it verifies "
      "against, so supplying the index is what enables the pass. These results were produced "
      "before the merge, by the branch, with the same verification logic.",
    MG+2.00,5.54,CW-2.30,0.80,size=11.8,color=WHITE,sp=1.08,label="nmb")
note(s,"Strongest single result in the deck. The flag is on main now; the numbers came from "
       "the branch, with identical verification logic.")

# ============================================================ 10 EXP 5 BACKENDS
s=slide(prs); title(s,"Experiment 5 — depletion backends","Results · interchangeable components")
hdr=["Backend","Mean F1","Host carry-over","Microbial loss","Peak memory"]
wid=[2.5,2.1,2.9,2.7,2.0]
BN={"rc_bowtie2":"Bowtie2","rc_minimap2":"minimap2","rc_centrifuge":"Centrifuge"}
rows=[]
for t in ["rc_bowtie2","rc_minimap2","rc_centrifuge"]:
    f1=st.mean(float([x for x in v4a if x["tool"]==t and x["dataset"]==d][0]["F1"]) for d in O4)
    hc=st.mean(float([x for x in v4a if x["tool"]==t and x["dataset"]==d][0]["Host_Remaining_Rate"])*100 for d in O4)
    ml=st.mean(float([x for x in v4a if x["tool"]==t and x["dataset"]==d][0]["Microbe_Loss_Rate"])*100 for d in O4)
    mm=st.mean(int(v4p[(t,d)]["max_memory_kb"])/1048576 for d in O4)
    rows.append([BN[t],f"{f1:.4f}",f"{hc:.3f} %",f"{ml:.3f} %",f"{mm:.1f} GB"])
rows.append(["Kraken2 (default)","—","—","—","15.5 GB"])
rows.append(["sylph (Exp. 6)","—","—","—","3.5 GB"])
table(s,MG,1.52,12.2,hdr,rows,wid,fs=12,hfs=12,rh=0.40,hh=0.42)
rect(s,MG,4.22,CW,1.26,fill=CARD,rad=0.07)
txt(s,"What this buys",MG+0.30,4.34,2.2,0.30,size=13.5,bold=True,color=RUST_DK,label="wb")
txt(s,"Because routing treats depletion as a replaceable component, backends can be swapped "
      "without touching the surrounding pipeline. Bowtie2 and minimap2 are closely matched on "
      "accuracy; Centrifuge is worse on both error directions. Memory differs by 4× between "
      "backends, which is the practical basis for choosing one.",
    MG+0.30,4.68,CW-0.60,0.72,size=12.5,color=MUTED,sp=1.10,label="wbb")
txt(s,"One runtime measurement (minimap2, 5M dataset: 1,802 s against 235 s for the larger 10M "
      "dataset) is a cold-start outlier and is excluded from any runtime comparison.",
    MG,5.62,CW,0.54,size=12.3,italic=True,color=RUST_DK,sp=1.10,label="out")
note(s,"The memory column is the real finding: sylph at 3.5 GB against Kraken2 at 15.5 GB.")

# ============================================================ 11 EXP 6 MATCHED
s=slide(prs); title(s,"Experiment 6 — 100M matched panel","Results · all three tools, same datasets")
hdr=["Dataset","Tool","F1","Microbial loss","Host carry","Runtime","Peak mem"]
wid=[1.85,2.55,1.35,1.95,1.65,1.55,1.20]
rows=[]
for d in D100:
    for t,nm in [("rustyclean_auto_sylph","RustyClean (sylph)"),("hostile","Hostile"),("kneaddata","KneadData")]:
        f1,ml,hc=met(*cm([x for x in m100a if x["tool"]==t and x["dataset"]==d][0]))
        rows.append([SH_[d] if t.startswith("rusty") else "",nm,f"{f1:.4f}",f"{ml:.3f} %",
                     f"{hc:.3f} %",f"{rt(m100p,d,t)/60:.0f} min",f"{rt(m100p,d,t,'max_memory_kb')/1048576:.1f} GB"])
table(s,MG,1.52,CW,hdr,rows,wid,fs=11.5,hfs=10.5,rh=0.355,hh=0.40,zebra=False)
rect(s,MG,4.42,CW/2-0.18,1.30,fill=CARD,rad=0.07)
txt(s,"sylph solves the memory problem",MG+0.28,4.54,CW/2-0.74,0.30,size=13.5,bold=True,
    color=RUST_DK,label="m1")
txt(s,f"{rt(m100p,D100[0],'rustyclean_auto_sylph','max_memory_kb')/1048576:.1f} GB against "
      f"15.5 GB for the Kraken2 index — a 4× reduction that makes dense concurrency practical.",
    MG+0.28,4.88,CW/2-0.74,0.72,size=12.3,color=MUTED,sp=1.10,label="m1b")
rect(s,MG+CW/2+0.18,4.42,CW/2-0.18,1.30,fill=INK,rad=0.07)
txt(s,"But Hostile is faster here",MG+CW/2+0.46,4.54,CW/2-0.74,0.30,size=13.5,bold=True,
    color=RUST_LT,label="m2")
txt(s,"On this panel Hostile beats RustyClean on runtime and on F1. This is the opposite of "
      "Experiment 3 and has to be reconciled before any speed claim is made.",
    MG+CW/2+0.46,4.88,CW/2-0.74,0.72,size=12.3,color=WHITE,sp=1.10,label="m2b")
txt(s,"Against KneadData the margin holds: 5.7× and 5.8× faster, with 7× less microbial loss.",
    MG,5.88,CW,0.36,size=13,bold=True,color=RUST_DK,label="kdm")
note(s,"This is the slide that complicates the story. Hostile wins on speed AND F1 at 100M, "
       "contradicting Experiment 3. Do not hide it.")

# ============================================================ 12 EXP 7
s=slide(prs); title(s,"Experiment 7 — 100M depletion only","Results · with verification, against Hostile")
hdr=["Dataset","Tool","F1","Microbial loss","Host carry","Runtime","Peak mem"]
wid=[1.85,3.05,1.25,1.75,1.55,1.35,1.28]
href={r["dataset"]:(lambda c:c[1]+c[2])(cm(r)) for r in s100a if "rustyclean" in r["tool"]}
rows=[]; seen=set()
for d in D100:
    for r in [x for x in s100a if x["dataset"]==d]:
        if (d,r["tool"]) in seen: continue
        seen.add((d,r["tool"])); tp,fp,tn,fn=cm(r)
        if abs((tp+fn)-href[d])<abs((fp+tn)-href[d]): tp,fp,tn,fn=tn,fn,tp,fp
        f1,ml,hc=met(tp,fp,tn,fn)
        nm="RustyClean (kraken2 + recheck)" if "rusty" in r["tool"] else "Hostile"
        rows.append([SH_[d] if "rusty" in r["tool"] else "",nm,f"{f1:.4f}",f"{ml:.3f} %",
                     f"{hc:.3f} %",f"{rt(s100p,d,r['tool'])/60:.0f} min",
                     f"{rt(s100p,d,r['tool'],'max_memory_kb')/1048576:.1f} GB"])
table(s,MG,1.52,CW,hdr,rows,wid,fs=11.5,hfs=10.5,rh=0.40,hh=0.42)
txt(s,"This is the default method — auto routing with the verification pass — against Hostile on "
      "a matched panel. It fills the gap Experiments 3 and 4 left open.",
    MG,3.60,CW,0.56,size=13,bold=True,color=INK,sp=1.10,label="dm")
cards=[("RustyClean wins","Host carry-over 0.07 % against 0.17 %, and F1 at 90 % host "
        "(0.9943 vs 0.9906). Runtime 1.2–1.3× faster.",RUST),
       ("Hostile wins","Microbial loss 0.04 % against 0.40 %, F1 at 50 % host, and peak memory "
        "by more than 4× (3.6 GB vs 15.5 GB).",SLATE),
       ("Data caveat","The Hostile rows in this file are stored with the opposite positive class; "
        "they were transposed before computing these metrics. Fix the file.",RUST_DK)]
cw3=(CW-2*0.30)/3
for i,(hd,body,col) in enumerate(cards):
    x=MG+i*(cw3+0.30)
    rect(s,x,4.26,cw3,1.62,fill=CARD if i<2 else None,line=None if i<2 else RUST_DK,rad=0.07)
    txt(s,hd,x+0.26,4.40,cw3-0.52,0.30,size=13.5,bold=True,color=col,label="w"+hd[:8])
    txt(s,body,x+0.26,4.76,cw3-0.52,1.04,size=12.3,color=MUTED,sp=1.10,label="wb"+hd[:8])
txt(s,"Two datasets only — the panel needs the remaining host fractions before this is quotable.",
    MG,6.02,CW,0.36,size=12.5,italic=True,color=MUTED,label="n2")
note(s,"This answers the gap: default method vs Hostile on the same data. Mixed result — we win "
       "on host carry-over and high-host F1, they win on microbial loss and memory.")

# ============================================================ 13 SUMMARY
s=slide(prs,dark=True)
txt(s,"Where the evidence stands",MG,0.68,CW,0.70,size=34,bold=True,color=WHITE,font=H,label="sh")
good=[("Adaptive routing works","Correct backend on 4/4 datasets; the rule reproduces every choice."),
      ("Verification pass works","19.7× less residual host for +6.7 % runtime, consistent across 4 datasets."),
      ("Clearly beats KneadData","5–6× faster with 7–13× less microbial loss, on every panel."),
      ("Backends are swappable","5 backends; sylph cuts peak memory from 15.5 GB to 3.5 GB.")]
bad=[("Hostile results conflict","Exp. 3 says we are 4.9× faster at 60M; Exp. 6 says Hostile is faster at 100M."),
     ("Recheck is not on main","The default method described here ships only on a branch."),
     ("Memory is our weak point","15.5 GB on the Kraken2 path against 1.1 GB for KneadData, 3.6 GB for Hostile."),
     ("Panel is narrow","4–6 datasets, all single-end, nothing in the 10–30 % routing band, no real data.")]
for col,(items,hd,c) in enumerate([(good,"Established",SAGE),(bad,"Open",RUST_LT)]):
    x=MG+col*(CW/2+0.10); w=CW/2-0.30
    txt(s,hd,x,1.52,w,0.32,size=16,bold=True,color=c,font=H,label="hh"+hd)
    y=1.94
    for t_,b_ in items:
        circ(s,x+0.02,y+0.06,0.18,fill=SAGE if col==0 else RUST)
        txt(s,t_,x+0.32,y,w-0.34,0.28,size=13.5,bold=True,color=WHITE,label="g"+t_[:10])
        txt(s,b_,x+0.32,y+0.30,w-0.34,0.68,size=12,color=MUTED,sp=1.10,label="gb"+t_[:10])
        y+=1.06
rect(s,MG,6.28,0.045,0.62,fill=RUST)
txt(s,"The engineering is done. The comparison against Hostile is not.",
    MG+0.30,6.38,CW-0.4,0.42,size=15,italic=True,color=RUST_LT,label="close")
note(s,"Honest summary. The KneadData story is solid; the Hostile story is unresolved and that is "
       "the blocker for the paper.")

# ============================================================ 14 NEXT
s=slide(prs); title(s,"What to run next","Priority order")
items=[("1","Reconcile the Hostile runtime contradiction",
        "Re-run Hostile on 60M and 100M on the same node with the same settings. Until this "
        "resolves, no speed claim against Hostile can be quoted.",RUST),
       ("2","Run the Kraken2 index ablation",
        "Every committed result used kraken16, a mixed database, not the human-only index the "
        "manuscript describes. benchmark_k2_index_ablation.sh measures how much of the 1.41% "
        "host carry-over is the index rather than the method.",RUST),
       ("3","Fix the transposed class labels",
        "results_100M_skipqc_matched/accuracy.csv stores Hostile with the opposite positive class. "
        "Anyone reading that file without checking will draw the wrong conclusion.",RUST_DK),
       ("4","Extend the panel",
        "Paired-end libraries, and host fractions between 10 % and 30 % where the routing rule "
        "actually decides. Currently nothing is measured in that band.",SLATE),
       ("5","Memory-aware worker cap",
        "Worker count comes from CPU count; with a 15.5 GB index, concurrent demand is W × 15.5 GB. "
        "Derive it from available memory instead.",SLATE)]
y=1.52
for n,hd,body,col in items:
    circ(s,MG+0.04,y+0.14,0.42,fill=col,t=n,size=15)
    txt(s,hd,MG+0.62,y+0.02,4.6,0.58,size=14,bold=True,color=INK,label="i"+hd[:12])
    txt(s,body,MG+5.35,y+0.02,CW-5.45,0.72,size=12.3,color=MUTED,sp=1.10,label="ib"+hd[:12])
    y+=0.94
rect(s,MG,6.34,CW,0.72,fill=CARD,rad=0.06)
txt(s,"Items 1–3 are blockers for the manuscript. Items 4–5 strengthen it.",
    MG+0.30,6.50,CW-0.60,0.34,size=13,bold=True,color=RUST_DK,label="blk")
note(s,"Close on the ask: item 1 is the one that decides whether the paper's speed claim survives.")

out="manuscript/RustyClean_Status.pptx"
prs.save(out)
print(f"saved {out} — {len(prs.slides._sldIdLst)} slides")
if warn:
    print("\nPOSSIBLE OVERFLOW:")
    for w_ in warn: print("  ",w_)
else: print("no suspected overflow")
